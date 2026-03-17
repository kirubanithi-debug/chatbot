"""
Generate lowchat.pptx — Law Chatbot Project Presentation
RICH colourful design with gradient-style headers, decorative shapes,
vibrant accent colours, and modern professional styling.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Rich Colour Palette ─────────────────────────────────────────────────────
WHITE         = RGBColor(255, 255, 255)
BLACK         = RGBColor(0, 0, 0)
OFF_WHITE     = RGBColor(248, 250, 252)

# Primary colours (black / teal)
NAVY          = RGBColor(10, 10, 10)
DARK_BLUE     = RGBColor(30, 30, 30)
TEAL          = RGBColor(0, 128, 128)
LIGHT_TEAL    = RGBColor(0, 180, 170)

# Accent colours
GOLD          = RGBColor(255, 193, 7)
CORAL         = RGBColor(255, 107, 107)
MINT          = RGBColor(100, 220, 180)
SKY_BLUE      = RGBColor(90, 90, 90)
SOFT_PURPLE   = RGBColor(156, 136, 255)
WARM_ORANGE   = RGBColor(255, 152, 0)

# Content colours
BODY_TEXT      = RGBColor(30, 30, 30)
SUBTITLE_TEXT  = RGBColor(80, 80, 80)
LIGHT_BG       = RGBColor(240, 240, 240)
CARD_BG        = RGBColor(245, 245, 245)
CARD_BORDER    = RGBColor(200, 200, 200)
HEADER_BG      = RGBColor(15, 15, 15)
FOOTER_BG_DARK = RGBColor(10, 10, 10)

# APRM card colours
BLUE_CARD      = RGBColor(220, 220, 220)
GREEN_CARD     = RGBColor(200, 240, 210)
ORANGE_CARD    = RGBColor(255, 230, 200)
PURPLE_CARD    = RGBColor(225, 215, 250)
BLUE_CARD_BD   = RGBColor(80, 80, 80)
GREEN_CARD_BD  = RGBColor(60, 160, 100)
ORANGE_CARD_BD = RGBColor(220, 140, 50)
PURPLE_CARD_BD = RGBColor(130, 100, 200)

# ── Fixed positions ─────────────────────────────────────────────────────────
SLIDE_W    = 10.0
MARGIN_L   = 0.55
MARGIN_R   = 0.55
TITLE_TOP  = 0.10
TITLE_H    = 0.80
BODY_TOP   = 1.30
BODY_H     = 5.50
CONTENT_W  = SLIDE_W - MARGIN_L - MARGIN_R


# ── Helpers ──────────────────────────────────────────────────────────────────

def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OFF_WHITE
    return slide


def _add_header_band(slide):
    """Wide dark-blue header band across the top."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(1.05),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = HEADER_BG
    sh.line.fill.background()


def _add_accent_dot(slide):
    """Small gold circle as a decorative brand dot in the header."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.35), Inches(0.30), Inches(0.40), Inches(0.40),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = GOLD
    sh.line.fill.background()


def _add_side_stripe(slide):
    """Thin teal stripe on the left."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.05), Inches(0.06), Inches(6.45),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = TEAL
    sh.line.fill.background()


def _add_footer_bar(slide, num):
    """Dark footer bar with gold slide number."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.10), Inches(SLIDE_W), Inches(0.40),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = FOOTER_BG_DARK
    sh.line.fill.background()
    # Gold accent line
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.08), Inches(SLIDE_W), Inches(0.03),
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = GOLD
    ln.line.fill.background()
    # Project name
    fb = slide.shapes.add_textbox(Inches(0.35), Inches(7.11), Inches(5), Inches(0.38))
    p = fb.text_frame.paragraphs[0]
    p.text = "Law Chatbot \u2014 AI Legal Guidance"
    _font(p, size=8, color=RGBColor(180, 195, 220))
    # Slide number in gold circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.25), Inches(7.12), Inches(0.32), Inches(0.32),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOLD
    circ.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(9.25), Inches(7.13), Inches(0.32), Inches(0.32))
    p = nb.text_frame.paragraphs[0]
    p.text = str(num)
    _font(p, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def _add_bottom_decor(slide):
    """Decorative teal triangle at bottom-right."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(8.8), Inches(6.2), Inches(1.2), Inches(0.90),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = LIGHT_TEAL
    sh.line.fill.background()


def _font(para, size=18, bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT, name="Calibri"):
    para.font.name = name
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = align


def _add_title(slide, text, align=PP_ALIGN.LEFT, size=26):
    """Title text rendered on top of the dark header band."""
    box = slide.shapes.add_textbox(
        Inches(0.90), Inches(TITLE_TOP + 0.15), Inches(CONTENT_W - 0.40), Inches(TITLE_H),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _font(p, size=size, bold=True, align=align, color=WHITE)
    return box


def _add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _body_box(slide, top=BODY_TOP, height=BODY_H):
    box = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(top), Inches(CONTENT_W - 0.30), Inches(height),
    )
    box.text_frame.word_wrap = True
    return box


def _add_content_card(slide, left, top, width, height):
    """White rounded card with subtle border."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = CARD_BORDER
    sh.line.width = Pt(1.2)
    return sh


def _write_bullets(tf, bullets, size=17, spacing=7, color=BODY_TEXT, bullet_char="\u25b8"):
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {text}" if bullet_char else text
        p.space_after = Pt(spacing)
        _font(p, size=size, color=color)


def _decorated_slide(slide, num):
    """Apply common decorations to a content slide."""
    _add_header_band(slide)
    _add_accent_dot(slide)
    _add_side_stripe(slide)
    _add_bottom_decor(slide)
    if num:
        _add_footer_bar(slide, num)


def _bullet_slide(prs, title, bullets, notes="", size=16, num=0):
    slide = _blank_slide(prs)
    _decorated_slide(slide, num)
    _add_title(slide, title)
    _add_content_card(slide, MARGIN_L, 1.20, CONTENT_W, 5.70)
    box = _body_box(slide, top=1.35, height=5.40)
    _write_bullets(box.text_frame, bullets, size=size)
    if notes:
        _add_notes(slide, notes)
    return slide


def _module_slide(prs, slide_title, mod_name, paragraphs, notes="", num=0):
    slide = _blank_slide(prs)
    _decorated_slide(slide, num)
    _add_title(slide, slide_title)

    # Module sub-header with teal accent
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_L), Inches(1.15), Inches(CONTENT_W), Inches(0.42),
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = TEAL
    hdr_bg.line.fill.background()

    hdr = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(1.18), Inches(CONTENT_W - 0.30), Inches(0.38),
    )
    p = hdr.text_frame.paragraphs[0]
    p.text = mod_name
    _font(p, size=15, bold=True, color=WHITE)

    _add_content_card(slide, MARGIN_L, 1.65, CONTENT_W, 5.25)

    box = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(1.75), Inches(CONTENT_W - 0.30), Inches(5.05),
    )
    box.text_frame.word_wrap = True
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(7)
        _font(p, size=12, color=BODY_TEXT)

    if notes:
        _add_notes(slide, notes)
    return slide


def _screenshot_slide(prs, title, caption, notes="", num=0):
    slide = _blank_slide(prs)
    _decorated_slide(slide, num)
    _add_title(slide, title)

    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(1.50), Inches(8.0), Inches(4.30),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_BG
    rect.line.color.rgb = TEAL
    rect.line.width = Pt(2)
    rect.line.dash_style = 4
    tf = rect.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "[ Insert Screenshot Here ]"
    _font(p, size=22, color=TEAL, align=PP_ALIGN.CENTER)

    cap = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(0.45))
    p = cap.text_frame.paragraphs[0]
    p.text = caption
    _font(p, size=12, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

    if notes:
        _add_notes(slide, notes)
    return slide


# ========================== SLIDE BUILDERS ==================================

def s01_title(prs):
    slide = _blank_slide(prs)

    # Full dark background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Top gold bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.10),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()

    # Decorative circle top-right
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.5), Inches(-1.5), Inches(5.0), Inches(5.0),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = DARK_BLUE
    circ.line.fill.background()

    # Decorative circle bottom-left
    circ2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.5), Inches(5.0), Inches(4.0), Inches(4.0),
    )
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = DARK_BLUE
    circ2.line.fill.background()

    # Gold accent bar top
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.2), Inches(2.0), Inches(3.6), Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Main title
    t = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(2.30), Inches(CONTENT_W), Inches(1.2),
    )
    p = t.text_frame.paragraphs[0]
    p.text = "Law Chatbot Project"
    _font(p, size=44, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

    # Subtitle
    s = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(3.50), Inches(CONTENT_W), Inches(0.6),
    )
    p = s.text_frame.paragraphs[0]
    p.text = "AI Based Indian Legal Guidance System"
    _font(p, size=22, bold=False, align=PP_ALIGN.CENTER, color=GOLD)

    # Gold accent bar bottom
    accent2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.2), Inches(4.30), Inches(3.6), Inches(0.06),
    )
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = GOLD
    accent2.line.fill.background()

    # Student details
    d = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(4.60), Inches(CONTENT_W), Inches(2.5),
    )
    tf = d.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = [
        ("Student Name:  [YOUR NAME]", 19, False),
        ("Department:  [YOUR DEPARTMENT]", 17, False),
        ("College:  [YOUR COLLEGE]", 17, False),
        ("Guide:  [GUIDE NAME]", 17, False),
    ]
    for i, (txt, sz, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.space_after = Pt(8)
        _font(p, size=sz, bold=bld, align=PP_ALIGN.CENTER, color=RGBColor(190, 210, 235))

    # Bottom gold bar
    bot_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.40), Inches(SLIDE_W), Inches(0.10),
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLD
    bot_bar.line.fill.background()

    _add_notes(slide, "Introduce the project title, your name, department, college, and guide name.")


def s02_abstract(prs):
    slide = _blank_slide(prs)
    _decorated_slide(slide, 2)
    _add_title(slide, "Abstract")
    _add_content_card(slide, MARGIN_L, 1.20, CONTENT_W, 5.70)
    box = _body_box(slide, top=1.35, height=5.40)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    paras = [
        "In India, many common people face legal problems every day but they do not know "
        "which law or section applies to their situation. Visiting a lawyer for every small "
        "question is costly and takes a lot of time. There is a clear need for a simple tool "
        "that can give basic legal direction to anyone instantly.",

        "This project builds an AI-powered legal chatbot where users can type their legal "
        "issue in plain everyday language through a clean web-based chat interface. The system "
        "uses a Retrieval Augmented Generation (RAG) approach to find the most relevant Indian "
        "law sections from a curated knowledge base stored in ChromaDB vector database.",

        "A dedicated Query Understanding module first cleans the user input by fixing spelling "
        "mistakes, expanding colloquial phrases, and detecting the legal intent of the question. "
        "This ensures that even poorly typed or informal queries are matched accurately.",

        "The backend is built with Django REST Framework and uses Celery with Redis for "
        "asynchronous processing. This means the React frontend stays fast and responsive "
        "while the AI computation runs in the background without any delay to the user.",

        "Once the relevant legal context is retrieved, the Ollama LLM generates a structured "
        "response in a fixed format containing LAW, SECTION, PUNISHMENT, NEXT STEPS, and "
        "DISCLAIMER. This gives users a clear and consistent answer every time they ask.",

        "The frontend also provides 55+ quick-question buttons so users can ask common legal "
        "queries with a single click. Every response includes a safety disclaimer advising "
        "users to consult a qualified lawyer for final decisions. The overall goal is to bridge "
        "the legal awareness gap for ordinary citizens across India.",
    ]
    for i, text in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        _font(p, size=12, color=BODY_TEXT)
    _add_notes(slide, "Abstract: problem, solution, tech approach, outcome.")


def s03_existing(prs):
    _bullet_slide(prs, "Existing System", [
        "Users currently search for legal information on Google, government websites, or legal forums which give unstructured and often confusing results.",
        "People visit advocates or legal-aid centres even for basic questions, which costs money and takes time.",
        "Most legal websites use complex legal language that ordinary people cannot easily understand.",
        "There is no single system where a user can type a simple question and get a clear answer with law, section, punishment, and next steps.",
        "Existing search tools do not handle spelling mistakes or informal language, so users must know the exact legal terms to find anything useful.",
    ], size=14, num=3,
       notes="Explain how legal guidance works today: manual search, lawyer visit, helpline.")


def s04_disadvantage(prs):
    _bullet_slide(prs, "Disadvantages of Existing System", [
        "Time-consuming \u2014 Users spend hours searching across multiple websites without finding a clear answer.",
        "Expensive \u2014 Even a basic query like 'what to do if my bike is stolen' may need a paid lawyer visit.",
        "Hard to understand \u2014 Legal documents use formal language that ordinary people find very difficult to read.",
        "No typo handling \u2014 Search engines fail if the user types 'sallary not payed' or 'harrasment'.",
        "No structured answers \u2014 Results are long paragraphs; users must manually find the relevant law and section.",
        "Not available 24/7 \u2014 Helplines and legal offices only work during business hours.",
    ], size=14, num=4,
       notes="Highlight each disadvantage clearly. These justify why the proposed system is needed.")


def s05_proposed(prs):
    _bullet_slide(prs, "Proposed System", [
        "An AI-powered legal chatbot where users type their issue in plain language and get a structured legal answer instantly.",
        "Built with React frontend, Django REST backend, Celery async workers, and Ollama LLM running locally.",
        "The Query Understanding module fixes typos, expands informal phrases, and detects the legal intent before searching.",
        "A RAG pipeline searches the ChromaDB vector database for the most relevant Indian law sections using semantic similarity.",
        "The Ollama LLM generates a response in a fixed format: LAW, SECTION, PUNISHMENT, NEXT STEPS, and DISCLAIMER.",
        "The frontend includes 55+ quick-question buttons, real-time polling, and a mandatory safety disclaimer on every response.",
    ], size=14, num=5,
       notes="Walk through the flow: user query -> query understanding -> RAG retrieval -> LLM -> structured response.")


def s06_advantage(prs):
    _bullet_slide(prs, "Advantages of Proposed System", [
        "Instant response \u2014 Users get a clear legal answer within seconds.",
        "Handles typos and informal language \u2014 No need to type perfect legal terms.",
        "Structured output \u2014 Every response follows LAW / SECTION / PUNISHMENT / NEXT STEPS / DISCLAIMER format.",
        "Available 24/7 \u2014 Works anytime as a web application, no office hours needed.",
        "Safe and responsible \u2014 Every answer includes a disclaimer to consult a real lawyer.",
        "Privacy friendly \u2014 AI runs locally using Ollama, no data sent to external cloud services.",
    ], size=14, num=6,
       notes="Keep advantages simple and direct.")


def s07_architecture(prs):
    slide = _blank_slide(prs)
    _decorated_slide(slide, 7)
    _add_title(slide, "APRM Architecture Diagram")

    # 4 APRM boxes with rich colours
    row_y = Inches(1.25)
    bw = Inches(2.10)
    bh = Inches(1.55)
    gap = Inches(0.12)
    x0 = Inches(MARGIN_L)

    card_fills = [BLUE_CARD, GREEN_CARD, ORANGE_CARD, PURPLE_CARD]
    card_borders = [BLUE_CARD_BD, GREEN_CARD_BD, ORANGE_CARD_BD, PURPLE_CARD_BD]
    data = [
        ("A \u2014 Analyze",   "Query Understanding\nTypo Fix\nIntent Detection\nPhrase Expansion"),
        ("P \u2014 Process",   "Django REST API\nCelery Task Queue\nRedis Broker\nAsync Dispatch"),
        ("R \u2014 Retrieve",  "ChromaDB Vectors\nKeyword Fallback\nSynonym Expansion\nConfidence Score"),
        ("M \u2014 Model",     "Ollama LLM\nStrict Prompt\nFallback Chain\nStructured Output"),
    ]

    boxes = []
    for i, (head, body) in enumerate(data):
        left = x0 + i * (bw + gap)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, row_y, bw, bh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = card_fills[i]
        sh.line.color.rgb = card_borders[i]
        sh.line.width = Pt(2)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = head
        _font(p, size=13, bold=True, align=PP_ALIGN.CENTER, color=NAVY)
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            p2.text = line
            _font(p2, size=9, color=BODY_TEXT, align=PP_ALIGN.CENTER)
        boxes.append(sh)

    # Arrow connectors (teal)
    for i in range(3):
        x1 = boxes[i].left + boxes[i].width
        y  = boxes[i].top + boxes[i].height // 2
        x2 = boxes[i + 1].left
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
        c.line.color.rgb = TEAL
        c.line.width = Pt(2.5)

    # Technology boxes (colourful pills)
    my = Inches(3.15)
    mw = Inches(1.60)
    mh = Inches(0.55)
    mg = Inches(0.18)
    mx0 = Inches(MARGIN_L + 0.15)

    tech_colors = [CORAL, SKY_BLUE, MINT, WARM_ORANGE, SOFT_PURPLE]
    comps = ["React", "SQLite", "ChromaDB", "Redis", "Ollama"]
    for i, lbl in enumerate(comps):
        left = mx0 + i * (mw + mg)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, my, mw, mh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = tech_colors[i]
        sh.line.fill.background()
        p = sh.text_frame.paragraphs[0]
        p.text = lbl
        _font(p, size=10, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

    # Data flow in a card
    _add_content_card(slide, MARGIN_L, 3.95, CONTENT_W, 2.95)

    fb = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(4.05), Inches(CONTENT_W - 0.30), Inches(2.75),
    )
    tf = fb.text_frame
    tf.word_wrap = True
    tf.clear()
    flow = [
        ("Data Flow:", True, 12, NAVY),
        ("1.  User types query in React UI  \u2192  sent to Django REST API", False, 10, BODY_TEXT),
        ("2.  API saves record  \u2192  dispatches Celery task via Redis", False, 10, BODY_TEXT),
        ("3.  Worker runs Query Understanding: typo fix, intent detection", False, 10, BODY_TEXT),
        ("4.  Enriched query searches ChromaDB for matching legal entries", False, 10, BODY_TEXT),
        ("5.  High confidence  \u2192  direct answer; else context sent to Ollama LLM", False, 10, BODY_TEXT),
        ("6.  LLM generates: LAW \u00b7 SECTION \u00b7 PUNISHMENT \u00b7 NEXT STEPS \u00b7 DISCLAIMER", False, 10, BODY_TEXT),
        ("7.  Response saved  \u2192  frontend polls  \u2192  answer shown to user", False, 10, BODY_TEXT),
    ]
    for i, (txt, bld, sz, clr) in enumerate(flow):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.space_after = Pt(2)
        _font(p, size=sz, bold=bld, color=clr)

    _add_notes(slide, "Walk through: APRM coloured boxes, tech components, then data flow.")


def s08_modules_overview(prs):
    _bullet_slide(prs, "Modules Overview", [
        "Frontend Module \u2014 React chat interface with quick-question buttons.",
        "API Module \u2014 Django REST endpoints for session and query management.",
        "Async Processing Module \u2014 Celery + Redis for background AI computation.",
        "Query Understanding Module \u2014 Typo fix, intent detection, phrase expansion.",
        "RAG Retrieval Module \u2014 ChromaDB vector search with keyword fallback.",
        "LLM Response Module \u2014 Ollama model with strict prompt for structured answers.",
        "Database Module \u2014 SQLite for chat history, ChromaDB for vector embeddings.",
        "Safety and Fallback Module \u2014 Default guidance and mandatory legal disclaimer.",
    ], size=14, num=8,
       notes="List all eight modules briefly.")


def s09_mod_analyze(prs):
    _module_slide(prs, "APRM Module \u2014 Analyze Query",
        "A \u2014 Query Understanding Module", [
        "This module is the first step of the APRM pipeline. It takes the raw text that the user "
        "typed in the chat box and cleans it so that the rest of the system can understand it properly. "
        "Users often type with spelling mistakes and informal language, so this step is very important.",

        "First, the input is converted to lowercase and extra spaces or special characters are removed. "
        "Then a typo correction dictionary with 80+ common misspellings fixes words like 'sallary' to "
        "'salary' and 'harrasment' to 'harassment'. This makes the query much more accurate for searching.",

        "Next, a phrase alias system expands informal phrases into proper legal terms. For example, "
        "'no road in my street' becomes 'civic road infrastructure street grievance' so that the search "
        "engine can find related law sections more easily.",

        "The module also detects the legal intent by checking the query tokens against 30+ intent "
        "categories like theft, murder, harassment, cyber fraud, dowry, and property dispute. It uses "
        "unigram, bigram, and trigram matching to catch multi-word legal phrases.",

        "Finally, entities like location (home, bus stop, school), time (today, yesterday, night), and "
        "objects (bike, phone, salary) are extracted. All this information is combined into an enriched "
        "query string that is sent to the next stage of the pipeline.",
    ], notes="First stage of APRM. Emphasize typo handling and intent detection.", num=9)


def s10_mod_process(prs):
    _module_slide(prs, "APRM Module \u2014 Process Request",
        "P \u2014 API and Async Processing Module", [
        "This module handles the communication between the frontend and the AI backend. When a user "
        "submits a query, the React frontend sends an HTTP POST request to the Django REST API with "
        "the session ID and the query text.",

        "The API first creates a ChatSession record if one does not already exist, then creates a new "
        "LegalQuery record with status set to 'pending'. This record stores the user question and will "
        "later hold the AI-generated response once processing is complete.",

        "Instead of processing the query immediately (which would block the user), the API dispatches "
        "the task to a Celery worker using Redis as the message broker. The API returns HTTP 202 "
        "(Accepted) right away so the user sees a loading indicator while the work happens in the background.",

        "The Celery worker picks up the task and runs the full pipeline: query understanding, retrieval "
        "from ChromaDB, and LLM response generation. Once done, it saves the final answer to the database "
        "and marks the status as 'completed'.",

        "Meanwhile, the React frontend polls the backend every 2 seconds to check if the response is "
        "ready. As soon as the status changes to 'completed', the answer is displayed to the user. "
        "This async design keeps the UI fast and responsive at all times.",
    ], notes="P in APRM. Key point: async processing keeps the UI fast.", num=10)


def s11_mod_retrieve(prs):
    _module_slide(prs, "APRM Module \u2014 Retrieve Context",
        "R \u2014 RAG Retrieval Module", [
        "This module searches the legal knowledge base to find the most relevant Indian law entry "
        "for the user's question. The knowledge base is built from legal_sections.csv which contains "
        "columns: law, section, punishment, next_steps, and keywords covering IPC, IT Act, BNSS, and more.",

        "The primary search method is vector similarity. All legal entries are embedded into ChromaDB "
        "using the Ollama nomic-embed-text model. When a query comes in, similarity_search() finds "
        "the top matching documents based on how close their meaning is to the user question.",

        "A secondary scoring method called find_best_legal_entry() checks every record using a formula "
        "that combines token overlap (70 percent weight) with fuzzy string similarity (30 percent weight) "
        "plus a bonus for exact section name matches. This catches cases that vector search might miss.",

        "The system also uses 40+ synonym groups to expand query words automatically. For example, if the "
        "user says 'stolen', the system also searches for 'theft, robbery, snatch, pickpocket, missing'. "
        "If Ollama embeddings are unavailable, a pure keyword-based fallback scorer takes over.",

        "If the best match score is 0.42 or above, the entry is returned directly as the final answer "
        "without needing the LLM at all. Otherwise, the top 3 matching documents are combined and sent "
        "as context to the Ollama LLM for generating the final response.",
    ], notes="R in APRM. Explain vector search, keyword fallback, 0.42 threshold.", num=11)


def s12_mod_model(prs):
    _module_slide(prs, "APRM Module \u2014 Model Response",
        "M \u2014 LLM Response Generation Module", [
        "This is the final stage of the APRM pipeline. It takes the user query along with the "
        "retrieved legal context and generates a clean, structured answer using the Ollama LLM "
        "(llama3.2:3b model running locally on the server).",

        "The build_prompt() function creates a carefully designed prompt that tells the LLM to act "
        "as an Indian Law Expert. It includes strict instructions to produce the output EXACTLY in "
        "the format: LAW, SECTION, PUNISHMENT, NEXT STEPS, and DISCLAIMER. This keeps every answer consistent.",

        "Before calling the LLM, a health check is done on the Ollama API to make sure it is running. "
        "The prompt is then sent to the /api/generate endpoint with stream set to False. If the LLM "
        "returns an empty response or the output does not contain the expected LAW: marker, the system "
        "automatically falls back to a pre-formatted answer.",

        "The fallback chain has three levels: first try format_response_from_entry() using the direct "
        "database match, then try format_response_from_context() using the retrieved context, and "
        "finally use get_default_guidance_response() which gives general legal direction to the user.",

        "For common simple questions like greetings, how to file FIR, or emergency helpline numbers, "
        "the Simple QA module answers instantly without calling the LLM at all. Every response always "
        "ends with a disclaimer advising the user to consult a qualified lawyer for final legal advice.",
    ], notes="M in APRM. Highlight strict format enforcement and the three-level fallback.", num=12)


def s13_mod_frontend(prs):
    _module_slide(prs, "Core Module \u2014 Frontend",
        "Frontend Interaction Module (React + Tailwind CSS)", [
        "The frontend is built using React.js and styled with Tailwind CSS. It provides a clean, "
        "simple chat interface where users can type their legal question in a text input box and "
        "receive AI-generated legal guidance in a readable format.",

        "When the page loads, a new ChatSession is automatically created by sending a POST request "
        "to the backend API. Users can either type their own question or click on any of the 55+ "
        "quick-question chip buttons to ask common legal queries with a single click.",

        "After submitting a query, the UI shows a loading indicator while the backend processes the "
        "request. The frontend polls the backend every 2 seconds to check if the AI response is ready. "
        "Once completed, the answer is displayed in a clear preformatted text block.",

        "The response is shown in the fixed format: LAW, SECTION, PUNISHMENT, NEXT STEPS, and DISCLAIMER. "
        "Connection errors and backend-unavailable states are handled gracefully with user-visible error "
        "messages. The single-page design keeps the experience intuitive for all users.",
    ], notes="Explain the React frontend: chat UI, quick chips, polling, and error handling.", num=13)


def s14_mod_db_safety(prs):
    _module_slide(prs, "Core Module \u2014 Database and Safety",
        "Database, Persistence, and Safety Fallback", [
        "The system uses two databases. SQLite stores chat sessions and legal query records (user question, "
        "AI response, status, timestamps). ChromaDB stores vector embeddings of all legal sections on "
        "disk so that similarity searches are fast and data persists across restarts.",

        "At startup, the legal_sections.csv file is loaded using Pandas. The columns are normalized, "
        "keywords are enriched using 50+ enrichment rule groups, and each entry is embedded into ChromaDB "
        "using the Ollama nomic-embed-text model. This builds the complete legal knowledge base.",

        "Each user conversation is grouped under a unique session ID, so the full question and answer "
        "history is maintained in order. This allows users to ask multiple questions in one session.",

        "The Safety Fallback module ensures that when retrieval confidence is low or no good match is found, "
        "the system returns general legal guidance instead of a wrong answer. A Simple QA module handles "
        "greetings and common FAQs (FIR filing, emergency numbers) without calling the full AI pipeline.",

        "Most importantly, every single response from the system \u2014 whether from the LLM, direct match, "
        "simple QA, or fallback \u2014 always includes a mandatory disclaimer telling the user to consult "
        "a qualified lawyer for final legal decisions.",
    ], notes="Cover database storage, data pipeline, and the safety/fallback system.", num=14)


def s15_tech_stack(prs):
    slide = _blank_slide(prs)
    _decorated_slide(slide, 15)
    _add_title(slide, "Technology Stack")

    tech_data = [
        ("Frontend",   "React.js, Vite, Tailwind CSS, Axios",            SKY_BLUE,     0),
        ("Backend",    "Django 5.1, Django REST Framework, Python 3.10",  TEAL,         1),
        ("Async",      "Celery (task queue), Redis (message broker)",     WARM_ORANGE,  2),
        ("AI / RAG",   "Ollama (llama3.2:3b), LangChain, ChromaDB",      SOFT_PURPLE,  3),
        ("Database",   "SQLite (relational), ChromaDB (vector storage)",  CORAL,        4),
        ("Data",       "Pandas, difflib (fuzzy matching)",                MINT,         5),
        ("Deploy",     "Docker and Docker Compose",                       GOLD,         6),
    ]

    start_y = 1.30
    card_h = 0.70
    card_gap = 0.08
    label_w = 1.80
    desc_w = CONTENT_W - label_w - 0.20

    for label, desc, color, idx in tech_data:
        y = start_y + idx * (card_h + card_gap)
        # Coloured label box
        lbl_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN_L), Inches(y), Inches(label_w), Inches(card_h),
        )
        lbl_sh.fill.solid()
        lbl_sh.fill.fore_color.rgb = color
        lbl_sh.line.fill.background()
        p = lbl_sh.text_frame.paragraphs[0]
        p.text = label
        _font(p, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # White description card
        desc_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN_L + label_w + 0.10), Inches(y), Inches(desc_w), Inches(card_h),
        )
        desc_sh.fill.solid()
        desc_sh.fill.fore_color.rgb = WHITE
        desc_sh.line.color.rgb = CARD_BORDER
        desc_sh.line.width = Pt(1)
        p = desc_sh.text_frame.paragraphs[0]
        p.text = desc
        _font(p, size=13, color=BODY_TEXT, align=PP_ALIGN.LEFT)

    _add_notes(slide, "List the complete tech stack.")


def s16_21_screenshots(prs):
    screens = [
        ("Screenshot \u2014 Chat Interface",   "React chat UI with quick-question chips and text input"),
        ("Screenshot \u2014 Query Pending",    "Loading state shown while Celery processes the query in background"),
        ("Screenshot \u2014 AI Response",      "Structured LAW / SECTION / PUNISHMENT / NEXT STEPS / DISCLAIMER output"),
        ("Screenshot \u2014 Quick Questions",   "One-click legal query chip buttons for 55+ common issues"),
        ("Screenshot \u2014 Architecture Flow", "Paste your system architecture diagram or flowchart image here"),
        ("Screenshot \u2014 Terminal / Docker", "Docker Compose running all services: backend, frontend, Celery, Redis, Ollama"),
    ]
    for i, (title, cap) in enumerate(screens):
        _screenshot_slide(prs, title, cap,
            notes=f"Replace the placeholder with an actual screenshot: {cap}",
            num=16 + i)


def s22_conclusion(prs):
    _bullet_slide(prs, "Conclusion", [
        "The Law Chatbot gives fast, first-level legal awareness to common users who lack legal knowledge.",
        "The APRM pipeline (Analyze \u2192 Process \u2192 Retrieve \u2192 Model) provides structured and reliable responses.",
        "Query understanding handles real-world challenges like typos, slang, and informal language.",
        "Multi-strategy retrieval (vector + keyword + fuzzy) improves accuracy over simple keyword search.",
        "Every response includes a safety disclaimer to encourage consulting a real lawyer.",
        "Future Scope: Multilingual support (Tamil, Hindi), voice input, mobile app, expanded legal data.",
    ], size=14, num=22,
       notes="Summarize project impact and future scope. End confidently.")


def s23_thank_you(prs):
    slide = _blank_slide(prs)

    # Full dark background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Top gold bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.10),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()

    # Decorative circles
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7.0), Inches(-1.0), Inches(4.5), Inches(4.5),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = DARK_BLUE
    circ.line.fill.background()

    circ2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.0), Inches(5.5), Inches(3.5), Inches(3.5),
    )
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = DARK_BLUE
    circ2.line.fill.background()

    # Centre accent card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.0), Inches(2.2), Inches(6.0), Inches(2.0),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BLUE
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)

    t = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.40), Inches(6.0), Inches(1.2),
    )
    p = t.text_frame.paragraphs[0]
    p.text = "Thank You"
    _font(p, size=46, bold=True, align=PP_ALIGN.CENTER, color=GOLD)

    s = slide.shapes.add_textbox(
        Inches(2.0), Inches(3.40), Inches(6.0), Inches(0.5),
    )
    p = s.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    _font(p, size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # Gold separator line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(4.60), Inches(3.0), Inches(0.05),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    d = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(5.0), Inches(CONTENT_W), Inches(0.5),
    )
    p = d.text_frame.paragraphs[0]
    p.text = "Law Chatbot \u2014 AI Based Indian Legal Guidance System"
    _font(p, size=14, color=RGBColor(180, 200, 230), align=PP_ALIGN.CENTER)

    # Bottom gold bar
    bot_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.40), Inches(SLIDE_W), Inches(0.10),
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLD
    bot_bar.line.fill.background()

    _add_notes(slide, "Thank the panel. Be ready for questions on architecture, modules, and demo.")


# ══════════════════════════════════════════════════════════════════════════════

def create_presentation(path: str):
    prs = Presentation()

    s01_title(prs)
    s02_abstract(prs)
    s03_existing(prs)
    s04_disadvantage(prs)
    s05_proposed(prs)
    s06_advantage(prs)
    s07_architecture(prs)
    s08_modules_overview(prs)
    s09_mod_analyze(prs)
    s10_mod_process(prs)
    s11_mod_retrieve(prs)
    s12_mod_model(prs)
    s13_mod_frontend(prs)
    s14_mod_db_safety(prs)
    s15_tech_stack(prs)
    s16_21_screenshots(prs)
    s22_conclusion(prs)
    s23_thank_you(prs)

    prs.save(path)
    return prs


if __name__ == "__main__":
    p = create_presentation("lowchat.pptx")
    print(f"Created: lowchat.pptx  ({len(p.slides)} slides)")
