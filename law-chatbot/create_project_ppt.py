"""
Generate law_chatbot_project_presentation.pptx
RICH colourful design: navy/gold/teal theme with decorative shapes,
content cards, coloured APRM boxes, and modern professional styling.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Rich Colour Palette ─────────────────────────────────────────────────────
WHITE         = RGBColor(255, 255, 255)
BLACK         = RGBColor(0, 0, 0)
OFF_WHITE     = RGBColor(248, 250, 252)

NAVY          = RGBColor(10, 10, 10)
DARK_BLUE     = RGBColor(30, 30, 30)
TEAL          = RGBColor(0, 128, 128)
LIGHT_TEAL    = RGBColor(0, 180, 170)

GOLD          = RGBColor(255, 193, 7)
CORAL         = RGBColor(255, 107, 107)
MINT          = RGBColor(100, 220, 180)
SKY_BLUE      = RGBColor(90, 90, 90)
SOFT_PURPLE   = RGBColor(156, 136, 255)
WARM_ORANGE   = RGBColor(255, 152, 0)

BODY_TEXT      = RGBColor(30, 30, 30)
SUBTITLE_TEXT  = RGBColor(80, 80, 80)
CARD_BORDER    = RGBColor(200, 200, 200)
HEADER_BG      = RGBColor(15, 15, 15)
FOOTER_BG_DARK = RGBColor(10, 10, 10)

BLUE_CARD      = RGBColor(220, 220, 220)
GREEN_CARD     = RGBColor(200, 240, 210)
ORANGE_CARD    = RGBColor(255, 230, 200)
PURPLE_CARD    = RGBColor(225, 215, 250)
BLUE_CARD_BD   = RGBColor(80, 80, 80)
GREEN_CARD_BD  = RGBColor(60, 160, 100)
ORANGE_CARD_BD = RGBColor(220, 140, 50)
PURPLE_CARD_BD = RGBColor(130, 100, 200)

SLIDE_W  = 10.0
MARGIN_L = 0.55
CONTENT_W = SLIDE_W - MARGIN_L - 0.55


# ── Helpers ──────────────────────────────────────────────────────────────────

def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OFF_WHITE
    return slide


def _font(para, size=18, bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT, name="Calibri"):
    para.font.name = name
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = align


def _add_header_band(slide):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(1.05),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = HEADER_BG
    sh.line.fill.background()


def _add_accent_dot(slide):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.35), Inches(0.30), Inches(0.40), Inches(0.40),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = GOLD
    sh.line.fill.background()


def _add_side_stripe(slide):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), Inches(0.06), Inches(6.45),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = TEAL
    sh.line.fill.background()


def _add_bottom_decor(slide):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.8), Inches(6.2), Inches(1.2), Inches(0.90),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = LIGHT_TEAL
    sh.line.fill.background()


def _add_footer_bar(slide, num):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.10), Inches(SLIDE_W), Inches(0.40),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = FOOTER_BG_DARK
    sh.line.fill.background()
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.08), Inches(SLIDE_W), Inches(0.03),
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = GOLD
    ln.line.fill.background()
    fb = slide.shapes.add_textbox(Inches(0.35), Inches(7.11), Inches(5), Inches(0.38))
    p = fb.text_frame.paragraphs[0]
    p.text = "Law Chatbot \u2014 AI Legal Guidance"
    _font(p, size=8, color=RGBColor(180, 195, 220))
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.25), Inches(7.12), Inches(0.32), Inches(0.32),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOLD
    circ.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(9.25), Inches(7.13), Inches(0.32), Inches(0.32))
    p = nb.text_frame.paragraphs[0]
    p.text = str(num)
    _font(p, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def _decorated_slide(slide, num):
    _add_header_band(slide)
    _add_accent_dot(slide)
    _add_side_stripe(slide)
    _add_bottom_decor(slide)
    if num:
        _add_footer_bar(slide, num)


def _add_title(slide, text, size=26):
    box = slide.shapes.add_textbox(
        Inches(0.90), Inches(0.25), Inches(CONTENT_W - 0.40), Inches(0.80),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _font(p, size=size, bold=True, align=PP_ALIGN.LEFT, color=WHITE)


def _add_content_card(slide, left, top, width, height):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = CARD_BORDER
    sh.line.width = Pt(1.2)
    return sh


def _bullet_slide(prs, title_text, bullets, bullet_size=16, num=0):
    slide = _blank_slide(prs)
    _decorated_slide(slide, num)
    _add_title(slide, title_text)
    _add_content_card(slide, MARGIN_L, 1.20, CONTENT_W, 5.70)

    box = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(1.35), Inches(CONTENT_W - 0.30), Inches(5.40),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25b8  {text}"
        p.space_after = Pt(7)
        _font(p, size=bullet_size)

    return slide


def add_architecture_slide(prs):
    slide = _blank_slide(prs)
    _decorated_slide(slide, 6)
    _add_title(slide, "APRM Architecture Diagram")

    row_y = Inches(1.25)
    bw = Inches(2.10)
    bh = Inches(1.40)
    gap = Inches(0.12)
    x0 = Inches(MARGIN_L)

    card_fills = [BLUE_CARD, GREEN_CARD, ORANGE_CARD, PURPLE_CARD]
    card_borders = [BLUE_CARD_BD, GREEN_CARD_BD, ORANGE_CARD_BD, PURPLE_CARD_BD]
    labels = [
        "A - Analyze Query",
        "P - Process Request",
        "R - Retrieve Context",
        "M - Model Response",
    ]
    sublabels = [
        "Intent + Typo Fix",
        "API + Celery Queue",
        "Chroma Vector Search",
        "Ollama Structured Reply",
    ]

    boxes = []
    for i, label in enumerate(labels):
        left = x0 + i * (bw + gap)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, row_y, bw, bh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = card_fills[i]
        sh.line.color.rgb = card_borders[i]
        sh.line.width = Pt(2)
        tf = sh.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        _font(p1, size=14, bold=True, align=PP_ALIGN.CENTER, color=NAVY)
        p2 = tf.add_paragraph()
        p2.text = sublabels[i]
        _font(p2, size=11, align=PP_ALIGN.CENTER, color=BODY_TEXT)
        boxes.append(sh)

    for i in range(len(boxes) - 1):
        x1 = boxes[i].left + boxes[i].width
        y1 = boxes[i].top + boxes[i].height // 2
        x2 = boxes[i + 1].left
        y2 = boxes[i + 1].top + boxes[i + 1].height // 2
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = TEAL
        c.line.width = Pt(2.5)

    _add_content_card(slide, MARGIN_L, 3.00, CONTENT_W, 1.60)
    flow = slide.shapes.add_textbox(Inches(MARGIN_L + 0.15), Inches(3.10), Inches(CONTENT_W - 0.30), Inches(1.40))
    tf = flow.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "User Query \u2192 Query Understanding \u2192 Legal Retrieval \u2192 LLM Generation \u2192 Final Legal Guidance"
    _font(p, size=16, bold=True, align=PP_ALIGN.CENTER, color=NAVY)
    p = tf.add_paragraph()
    p.text = "Supports async processing, structured response format, and fallback safety message."
    _font(p, size=13, align=PP_ALIGN.CENTER, color=SUBTITLE_TEXT)


def add_module_brief_slide(prs, title_text, module_title, points, num=0):
    slide = _blank_slide(prs)
    _decorated_slide(slide, num)
    _add_title(slide, title_text)

    # Module sub-header with teal accent
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_L), Inches(1.15), Inches(CONTENT_W), Inches(0.42),
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = TEAL
    hdr_bg.line.fill.background()

    hdr = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(1.18), Inches(CONTENT_W - 0.30), Inches(0.38),
    )
    p = hdr.text_frame.paragraphs[0]
    p.text = module_title
    _font(p, size=15, bold=True, color=WHITE)

    _add_content_card(slide, MARGIN_L, 1.65, CONTENT_W, 5.25)

    content_box = slide.shapes.add_textbox(
        Inches(MARGIN_L + 0.15), Inches(1.80), Inches(CONTENT_W - 0.30), Inches(4.95),
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25b8  {point}"
        p.space_after = Pt(7)
        _font(p, size=16)


def create_presentation(output_path):
    prs = Presentation()

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # Top gold bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.10),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()
    # Decorative circles
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(6.5), Inches(-1.5), Inches(5.0), Inches(5.0),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = DARK_BLUE
    circ.line.fill.background()
    circ2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.0), Inches(4.0), Inches(4.0),
    )
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = DARK_BLUE
    circ2.line.fill.background()
    # Gold accent bars
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(2.5), Inches(3.6), Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(2.80), Inches(CONTENT_W), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    p.text = "Law Chatbot Project"
    _font(p, size=44, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

    s = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(3.90), Inches(CONTENT_W), Inches(0.6))
    p = s.text_frame.paragraphs[0]
    p.text = "AI Based Legal Guidance System | Clean Presentation"
    _font(p, size=20, align=PP_ALIGN.CENTER, color=GOLD)

    accent2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(4.70), Inches(3.6), Inches(0.06),
    )
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = GOLD
    accent2.line.fill.background()
    # Bottom gold bar
    bot_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.40), Inches(SLIDE_W), Inches(0.10),
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLD
    bot_bar.line.fill.background()

    # ── Content slides ───────────────────────────────────────────────────────

    _bullet_slide(prs, "Abstract", [
        "This project builds an AI-powered legal chatbot for first-level legal support.",
        "Users can ask in simple language and get structured legal guidance instantly.",
        "The response format is fixed: LAW, SECTION, PUNISHMENT, NEXT STEPS, DISCLAIMER.",
        "The core goal is accessibility: legal awareness for common people without legal jargon.",
        "It combines query understanding, legal retrieval, and LLM generation in one flow.",
    ], bullet_size=18, num=2)

    _bullet_slide(prs, "Existing System", [
        "People usually depend on manual web search or direct lawyer consultation.",
        "Legal information is spread across many sources and hard to verify quickly.",
        "Most users do not know which law or section applies to their issue.",
        "There is no single simple interface for initial legal understanding.",
        "Common language queries and typo-heavy queries are not handled well.",
    ], bullet_size=18, num=3)

    _bullet_slide(prs, "Disadvantage", [
        "Time-consuming and confusing for non-legal users.",
        "No real-time structured legal direction in one response.",
        "Low support for colloquial terms and spelling mistakes.",
        "Manual process does not scale for high number of routine questions.",
        "Lack of a safe fallback message can mislead users in uncertain cases.",
    ], bullet_size=18, num=4)

    _bullet_slide(prs, "Proposed System", [
        "Frontend (React) collects user query and displays chat history.",
        "Backend (Django REST) creates sessions, stores messages, and serves APIs.",
        "Celery + Redis execute heavy AI tasks asynchronously.",
        "RAG layer retrieves relevant laws from legal_sections.csv via Chroma vector DB.",
        "Ollama model generates final readable legal response in strict output format.",
    ], bullet_size=18, num=5)

    _bullet_slide(prs, "Advantage", [
        "Fast and clear first-level legal guidance.",
        "Improved relevance through query preprocessing + retrieval.",
        "Consistent response structure helps users understand quickly.",
        "Async architecture supports multiple users without blocking.",
        "Fallback + disclaimer keeps output safe when confidence is low.",
    ], bullet_size=18, num=6)

    add_architecture_slide(prs)

    _bullet_slide(prs, "Modules Overview", [
        "1) Frontend Interaction Module",
        "2) API and Session Management Module",
        "3) Async Task Processing Module",
        "4) Query Understanding Module",
        "5) RAG Retrieval Module",
        "6) LLM Response Generation Module",
        "7) Database and Persistence Module",
        "8) Safety and Fallback Module",
    ], bullet_size=18, num=7)

    add_module_brief_slide(prs, "APRM Module Brief", "A - Analyze Query", [
        "Accepts raw user text from chat UI.",
        "Fixes common spelling mistakes and phrase variations.",
        "Detects issue intent using legal-domain keyword mapping.",
        "Builds an enriched query for better legal retrieval.",
    ], num=8)

    add_module_brief_slide(prs, "APRM Module Brief", "P - Process Request", [
        "Validates request through serializer and creates legal query record.",
        "Pushes task to Celery worker through Redis queue.",
        "Prevents frontend delay by running heavy work in background.",
        "Tracks status as pending and completed for polling updates.",
    ], num=9)

    add_module_brief_slide(prs, "APRM Module Brief", "R - Retrieve Context", [
        "Loads legal knowledge base from structured CSV data.",
        "Converts legal content into embeddings and stores in Chroma DB.",
        "Finds top relevant sections using semantic similarity search.",
        "Returns law, section, punishment, and next-step context.",
    ], num=10)

    add_module_brief_slide(prs, "APRM Module Brief", "M - Model Response", [
        "Sends context and user query to Ollama model with strict prompt.",
        "Generates readable legal answer in fixed structured format.",
        "Uses fallback formatter if model service is unavailable.",
        "Always includes safety disclaimer to consult a legal professional.",
    ], num=11)

    add_module_brief_slide(prs, "System Module Brief", "Frontend + API Modules", [
        "Frontend provides chat window, quick prompts, and status messages.",
        "API module exposes endpoints for session, submit-query, and messages.",
        "Polling every few seconds gives near real-time updates.",
        "Simple UI keeps interaction easy for non-technical users.",
    ], num=12)

    add_module_brief_slide(prs, "System Module Brief", "Database + Safety Modules", [
        "SQLite stores sessions, questions, responses, and timestamps.",
        "Chroma persistent storage keeps vector index for legal retrieval.",
        "Fallback module handles weak-context and no-context situations safely.",
        "Ensures system gives useful guidance without overconfident claims.",
    ], num=13)

    _bullet_slide(prs, "Conclusion", [
        "Law Chatbot provides practical and fast first-level legal awareness.",
        "The APRM pipeline makes responses structured, relevant, and readable.",
        "Modular architecture helps future scaling and feature extension.",
        "Future scope: multilingual support, voice input, and richer legal datasets.",
    ], bullet_size=18, num=14)

    # ── Thank You slide ──────────────────────────────────────────────────────
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.10),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.0), Inches(-1.0), Inches(4.5), Inches(4.5),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = DARK_BLUE
    circ.line.fill.background()
    circ2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.5), Inches(3.5), Inches(3.5),
    )
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = DARK_BLUE
    circ2.line.fill.background()
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.5), Inches(6.0), Inches(2.0),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BLUE
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)
    t = slide.shapes.add_textbox(Inches(2.0), Inches(2.70), Inches(6.0), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    p.text = "Thank You"
    _font(p, size=46, bold=True, align=PP_ALIGN.CENTER, color=GOLD)
    s = slide.shapes.add_textbox(Inches(2.0), Inches(3.70), Inches(6.0), Inches(0.5))
    p = s.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    _font(p, size=20, color=WHITE, align=PP_ALIGN.CENTER)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.60), Inches(3.0), Inches(0.05),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    bot_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.40), Inches(SLIDE_W), Inches(0.10),
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLD
    bot_bar.line.fill.background()

    prs.save(output_path)


if __name__ == "__main__":
    create_presentation("law_chatbot_project_presentation.pptx")
    print("Created: law_chatbot_project_presentation.pptx")
